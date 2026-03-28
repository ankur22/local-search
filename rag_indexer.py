#!/usr/bin/env python3
import argparse
import hashlib
import json
import os
import re
import sys
from contextlib import nullcontext
from dataclasses import dataclass
from typing import Iterable, List, Dict, Optional

import requests
from tqdm import tqdm

from pypdf import PdfReader
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
import markdown as md
from pptx import Presentation
from openpyxl import load_workbook

import chromadb
from chromadb.config import Settings
import unicodedata

try:
    from sigil_sdk import (
        Client as SigilClient,
        ClientConfig,
        EmbeddingResult,
        EmbeddingStart,
        GenerationExportConfig,
        ModelRef,
    )
    _SIGIL_AVAILABLE = True
except ImportError:
    _SIGIL_AVAILABLE = False

_sigil_client = None

# ------------------------ Config ------------------------

DEFAULT_EMBED_MODEL = "mxbai-embed-large"
DEFAULT_OLLAMA_URL = "http://localhost:11434"
DEFAULT_DB_DIR = "./chroma_db"
DEFAULT_COLLECTION = "my_corpus"

DEFAULT_CHUNK_CHARS = 1200
DEFAULT_OVERLAP_CHARS = 150

DEFAULT_CHUNK_TOKENS = 256
DEFAULT_OVERLAP_TOKENS = 64

MAX_FILE_BYTES_DEFAULT = 25 * 1024 * 1024  # 25 MB

PDF_EXTS = {".pdf"}
TEXT_EXTS = {".txt", ".md"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
XLSX_EXTS = {".xlsx"}
HTML_EXTS = {".html", ".htm"}
SOURCE_CODE_EXTS = {
    ".go", ".py", ".js", ".ts", ".tsx", ".java", ".kt", ".cs",
    ".rb", ".php", ".cpp", ".cc", ".c", ".h", ".hpp", ".rs",
    ".swift", ".sh", ".ps1", ".r", ".m", ".scala",
}

DEFAULT_INCLUDE_EXTS = list(
    PDF_EXTS
    | TEXT_EXTS
    | DOCX_EXTS
    | PPTX_EXTS
    | XLSX_EXTS
    | HTML_EXTS
    | SOURCE_CODE_EXTS
)

DEFAULT_EXCLUDE_DIRS = [
    ".obsidian", ".git", ".hg", ".svn",
    ".venv", "venv",
    "node_modules",
    "__pycache__",
    "dist", "build", "out", "target",
]

# tiktoken is optional: used for token-aware chunking if available
try:
    import tiktoken
    _ENC = tiktoken.get_encoding("cl100k_base")
except Exception:
    _ENC = None

# ------------------------ Types ------------------------

@dataclass
class Record:
    doc_id: str
    chunk_id: str
    text: str
    source_path: str
    ext: str
    chunk_index: int
    total_chunks: int

# ------------------------ File helpers ------------------------

def read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(path: str) -> str:
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            pass
    return "\n".join(parts).strip()

def read_docx(path: str) -> str:
    doc = DocxDocument(path)
    paras = [p.text for p in doc.paragraphs]
    return "\n".join(paras).strip()

def read_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        if slide.shapes and getattr(slide.shapes, "title", None) is not None:
            title = slide.shapes.title.text or ""
            if title:
                parts.append(title)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if hasattr(shape, "table"):
                try:
                    tbl = shape.table
                    for r in tbl.rows:
                        parts.append("\t".join([(c.text or "") for c in r.cells]))
                except Exception:
                    pass
    return "\n".join(parts).strip()

def read_xlsx(path: str) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [("" if v is None else str(v)) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()

def read_html(path: str) -> str:
    html = read_text_file(path)
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text("\n").strip()

def read_markdown(path: str) -> str:
    text = read_text_file(path)
    html = md.markdown(text)
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text("\n").strip()

def clean_text(s: str) -> str:
    # Normalise newlines
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    # Strip most control chars except newlines/tabs
    cleaned_chars = []
    for ch in s:
        cat = unicodedata.category(ch)
        if cat.startswith("C") and ch not in ("\n", "\t"):
            continue
        cleaned_chars.append(ch)
    s = "".join(cleaned_chars)
    # collapse >2 blank lines
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def extract_text_for_path(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in PDF_EXTS:
            txt = read_pdf(path)
        elif ext in DOCX_EXTS:
            txt = read_docx(path)
        elif ext in PPTX_EXTS:
            txt = read_pptx(path)
        elif ext in XLSX_EXTS:
            txt = read_xlsx(path)
        elif ext in TEXT_EXTS:
            if ext == ".md":
                txt = read_markdown(path)
            else:
                txt = read_text_file(path)
        elif ext in HTML_EXTS:
            txt = read_html(path)
        elif ext in SOURCE_CODE_EXTS:
            txt = read_text_file(path)
        else:
            return ""
        return clean_text(txt)
    except Exception as e:
        print(f"[WARN] Failed to parse {path}: {e}", file=sys.stderr)
        return ""

# ------------------------ Chunking ------------------------

def chunk_text_chars(s: str, size: int, overlap: int) -> List[str]:
    s = s.strip()
    if not s:
        return []
    out = []
    step = max(size - overlap, 1)
    i = 0
    while i < len(s):
        out.append(s[i:i + size])
        i += step
    return out

def chunk_text_tokens(s: str, target_tokens: int, overlap_tokens: int) -> List[str]:
    if not s:
        return []
    if _ENC is None:
        approx = target_tokens * 4
        approx_ov = overlap_tokens * 4
        return chunk_text_chars(s, approx, approx_ov)
    toks = _ENC.encode(s)
    out = []
    step = max(target_tokens - overlap_tokens, 1)
    i = 0
    while i < len(toks):
        chunk = toks[i:i + target_tokens]
        out.append(_ENC.decode(chunk))
        i += step
    return out

# ------------------------ Embeddings (Ollama) ------------------------

def ollama_embed(
    text: str,
    model: str = DEFAULT_EMBED_MODEL,
    base_url: str = DEFAULT_OLLAMA_URL,
    timeout: int = 60,
) -> List[float]:
    ctx = (
        _sigil_client.start_embedding(
            EmbeddingStart(
                agent_name="local-search-indexer",
                agent_version="0.1.0",
                model=ModelRef(provider="ollama", name=model),
            )
        )
        if _sigil_client is not None
        else nullcontext()
    )

    with ctx as rec:
        try:
            r = requests.post(
                f"{base_url}/api/embeddings",
                json={"model": model, "prompt": text},
                timeout=timeout,
            )
            r.raise_for_status()
            js = r.json()

            if "embedding" in js:
                embedding = js["embedding"]
            elif "data" in js and js["data"] and "embedding" in js["data"][0]:
                embedding = js["data"][0]["embedding"]
            else:
                raise RuntimeError(
                    f"Unexpected embeddings response from Ollama: {js}"
                )

            if rec is not None:
                rec.set_result(
                    EmbeddingResult(input_count=1, response_model=model)
                )

            return embedding
        except Exception as exc:
            if rec is not None:
                rec.set_call_error(exc)
            raise

# ------------------------ Discovery & IDs ------------------------

def discover_files(
    root: str,
    include_exts: Iterable[str],
    exclude_dirs: Iterable[str],
    max_bytes: int,
) -> List[str]:
    include_exts = {e.lower() for e in include_exts}
    exclude_dirs = {d.lower() for d in exclude_dirs}
    paths: List[str] = []

    for dirpath, dirnames, filenames in os.walk(root):
        base = os.path.basename(dirpath).lower()
        if base in exclude_dirs:
            continue
        dirnames[:] = [d for d in dirnames if d.lower() not in exclude_dirs]
        for fn in filenames:
            ext = os.path.splitext(fn)[1].lower()
            if ext in include_exts:
                full = os.path.join(dirpath, fn)
                try:
                    if os.path.getsize(full) <= max_bytes:
                        paths.append(full)
                except Exception:
                    pass
    return paths

def make_doc_id(path: str) -> str:
    # Stable ID: absolute path
    return os.path.abspath(path)

def file_signature(path: str) -> str:
    # Cheap change detector: size + mtime
    st = os.stat(path)
    key = f"{st.st_size}::{int(st.st_mtime)}"
    return hashlib.sha1(key.encode("utf-8")).hexdigest()

def make_chunk_id(doc_id: str, idx: int) -> str:
    return f"{doc_id}#chunk:{idx}"

# ------------------------ State (for incremental indexing) ------------------------

def get_state_path(db_dir: str) -> str:
    return os.path.join(db_dir, "index_state.json")

def load_state(db_dir: str) -> Dict[str, Dict]:
    path = get_state_path(db_dir)
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def save_state(db_dir: str, state: Dict[str, Dict]) -> None:
    os.makedirs(db_dir, exist_ok=True)
    path = get_state_path(db_dir)
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(state, f, indent=2)
    os.replace(tmp, path)

# ------------------------ Generic metadata helpers ------------------------

def infer_kind(source_path: str, ext: str) -> str:
    p = source_path.lower()
    if "release notes" in p or "release_notes" in p:
        return "release_notes"
    if ext in SOURCE_CODE_EXTS:
        return "code"
    if ext in {".md", ".txt"}:
        return "note"
    if ext in PDF_EXTS:
        return "pdf"
    return "other"

def extract_version(source_path: str, full_text: str) -> Optional[str]:
    # Try filename / path first
    m = re.search(r"v\d+\.\d+\.\d+", source_path)
    if m:
        return m.group(0)
    # Fallback: search in text
    m = re.search(r"v\d+\.\d+\.\d+", full_text)
    if m:
        return m.group(0)
    return None

# ------------------------ Chroma helpers ------------------------

def build_collection(db_dir: str, name: str):
    client = chromadb.PersistentClient(path=db_dir, settings=Settings(allow_reset=False))
    try:
        col = client.get_collection(name)
    except Exception:
        col = client.create_collection(name)
    return col

def refresh_document(
    collection,
    doc_id: str,
    source_path: str,
    text: str,
    ext: str,
    use_tokens: bool,
    chunk_chars: int,
    overlap_chars: int,
    chunk_tokens: int,
    overlap_tokens: int,
    embed_model: str,
    ollama_url: str,
):
    # Delete old chunks for this document
    collection.delete(where={"doc_id": doc_id})

    # Infer kind/version once per document
    kind = infer_kind(source_path, ext)
    version = extract_version(source_path, text)

    # Chunk
    if use_tokens:
        chunks = chunk_text_tokens(text, chunk_tokens, overlap_tokens)
    else:
        chunks = chunk_text_chars(text, chunk_chars, overlap_chars)

    total = len(chunks)
    if total == 0:
        return

    failed = 0
    embedded = 0

    for idx, ch in enumerate(chunks):
        chunk_id = make_chunk_id(doc_id, idx)
        try:
            emb = ollama_embed(ch, embed_model, ollama_url)
        except Exception as e:
            failed += 1
            snippet = ch[:200].replace("\n", "\\n")
            print(
                f"[ERROR] Embedding failed for {source_path} "
                f"(chunk {idx+1}/{total}, len={len(ch)}): {e}\n"
                f"         Snippet: {snippet!r}",
                file=sys.stderr,
            )
            continue

        metadata = {
            "doc_id": doc_id,
            "source": source_path,
            "ext": ext,
            "kind": kind,
            "version": version,
            "chunk_index": idx,
            "total_chunks": total,
        }
        collection.upsert(
            ids=[chunk_id],
            documents=[ch],
            metadatas=[metadata],
            embeddings=[emb],
        )
        embedded += 1

    if failed:
        print(
            f"[WARN] {failed} chunks failed to embed for {source_path}; "
            f"{embedded} chunks indexed successfully.",
            file=sys.stderr,
        )

# ------------------------ Main indexing logic ------------------------

def index_path(
    roots: List[str],
    db_dir: str,
    collection_name: str,
    include_exts: List[str],
    exclude_dirs: List[str],
    max_file_bytes: int,
    use_tokens: bool,
    chunk_chars: int,
    overlap_chars: int,
    chunk_tokens: int,
    overlap_tokens: int,
    embed_model: str,
    ollama_url: str,
    refresh_all: bool,
    prune_missing: bool,
):
    files: List[str] = []
    for root in roots:
        files.extend(
            discover_files(root, include_exts, exclude_dirs, max_file_bytes)
        )

    if not files:
        print("No files found for the given extensions.")
        return

    collection = build_collection(db_dir, collection_name)
    old_state = load_state(db_dir)
    new_state: Dict[str, Dict] = {}

    print(f"Discovered {len(files)} files. Indexing (refresh_all={refresh_all})...")

    for path in tqdm(files, desc="Indexing files"):
        abs_path = os.path.abspath(path)
        doc_id = make_doc_id(path)
        sig = file_signature(path)
        ext = os.path.splitext(path)[1].lower()

        new_state[doc_id] = {
            "sig": sig,
            "path": abs_path,
            "ext": ext,
        }

        if not refresh_all and doc_id in old_state and old_state[doc_id].get("sig") == sig:
            continue

        text = extract_text_for_path(path)
        if not text:
            continue

        refresh_document(
            collection=collection,
            doc_id=doc_id,
            source_path=abs_path,
            text=text,
            ext=ext,
            use_tokens=use_tokens,
            chunk_chars=chunk_chars,
            overlap_chars=overlap_chars,
            chunk_tokens=chunk_tokens,
            overlap_tokens=overlap_tokens,
            embed_model=embed_model,
            ollama_url=ollama_url,
        )

    if prune_missing:
        old_ids = set(old_state.keys())
        new_ids = set(new_state.keys())
        removed = old_ids - new_ids
        if removed:
            print(f"Pruning {len(removed)} documents that no longer exist on disk...")
            for doc_id in removed:
                collection.delete(where={"doc_id": doc_id})

    save_state(db_dir, new_state)
    print("Indexing complete.")

# ------------------------ CLI ------------------------

def main():
    ap = argparse.ArgumentParser(
        description="Index local documents into Chroma using Ollama embeddings."
    )
    ap.add_argument(
        "--root",
        action="append",
        required=True,
        help="Root directory to scan (can be used multiple times)",
    )
    ap.add_argument("--db", default=DEFAULT_DB_DIR, help="Chroma DB directory")
    ap.add_argument("--collection", default=DEFAULT_COLLECTION, help="Chroma collection name")
    ap.add_argument("--include-ext", nargs="*", default=DEFAULT_INCLUDE_EXTS,
                    help=f"File extensions to include (e.g. .pdf .md .docx). Default: {DEFAULT_INCLUDE_EXTS}")
    ap.add_argument("--exclude-dir", nargs="*", default=DEFAULT_EXCLUDE_DIRS,
                    help=f"Directory basenames to skip. Default: {DEFAULT_EXCLUDE_DIRS}")
    ap.add_argument("--max-file-bytes", type=int, default=MAX_FILE_BYTES_DEFAULT,
                    help="Skip files larger than this size (bytes)")

    ap.add_argument("--use-tokens", action="store_true",
                    help="Use token-aware chunking (requires tiktoken; falls back to chars otherwise)")
    ap.add_argument("--chunk-chars", type=int, default=DEFAULT_CHUNK_CHARS,
                    help="Chars per chunk (when not using --use-tokens)")
    ap.add_argument("--overlap-chars", type=int, default=DEFAULT_OVERLAP_CHARS,
                    help="Char overlap between chunks")
    ap.add_argument("--chunk-tokens", type=int, default=DEFAULT_CHUNK_TOKENS,
                    help="Tokens per chunk (with --use-tokens)")
    ap.add_argument("--overlap-tokens", type=int, default=DEFAULT_OVERLAP_TOKENS,
                    help="Token overlap between chunks (with --use-tokens)")

    ap.add_argument("--embed-model", default=DEFAULT_EMBED_MODEL,
                    help="Ollama embedding model name")
    ap.add_argument("--ollama-url", default=DEFAULT_OLLAMA_URL,
                    help="Base URL for Ollama (e.g. http://localhost:11434)")

    ap.add_argument("--refresh-all", action="store_true",
                    help="Reindex all files even if unchanged")
    ap.add_argument("--no-prune-missing", action="store_true",
                    help="Do not delete docs from the index that no longer exist on disk")

    args = ap.parse_args()

    global _sigil_client
    sigil_endpoint = os.environ.get("SIGIL_GENERATION_EXPORT_ENDPOINT", "")
    if _SIGIL_AVAILABLE and sigil_endpoint:
        otel_endpoint = os.environ.get("OTEL_EXPORTER_OTLP_ENDPOINT", "")
        if otel_endpoint:
            from opentelemetry import metrics, trace
            from opentelemetry.sdk.metrics import MeterProvider
            from opentelemetry.sdk.metrics.export import PeriodicExportingMetricReader
            from opentelemetry.sdk.trace import TracerProvider
            from opentelemetry.sdk.trace.export import BatchSpanProcessor
            from opentelemetry.exporter.otlp.proto.http.metric_exporter import OTLPMetricExporter
            from opentelemetry.exporter.otlp.proto.http.trace_exporter import OTLPSpanExporter
            from opentelemetry.sdk.resources import Resource

            resource = Resource.create({"service.name": "local-search-indexer"})
            tp = TracerProvider(resource=resource)
            tp.add_span_processor(BatchSpanProcessor(OTLPSpanExporter(endpoint=f"{otel_endpoint}/v1/traces")))
            trace.set_tracer_provider(tp)

            mp = MeterProvider(
                resource=resource,
                metric_readers=[PeriodicExportingMetricReader(
                    OTLPMetricExporter(endpoint=f"{otel_endpoint}/v1/metrics"),
                    export_interval_millis=5000,
                )],
            )
            metrics.set_meter_provider(mp)

        _sigil_client = SigilClient(
            ClientConfig(
                generation_export=GenerationExportConfig(
                    protocol="http",
                    endpoint=sigil_endpoint,
                ),
            )
        )

    try:
        index_path(
            roots=args.root,
            db_dir=args.db,
            collection_name=args.collection,
            include_exts=args.include_ext,
            exclude_dirs=args.exclude_dir,
            max_file_bytes=args.max_file_bytes,
            use_tokens=args.use_tokens,
            chunk_chars=args.chunk_chars,
            overlap_chars=args.overlap_chars,
            chunk_tokens=args.chunk_tokens,
            overlap_tokens=args.overlap_tokens,
            embed_model=args.embed_model,
            ollama_url=args.ollama_url,
            refresh_all=args.refresh_all,
            prune_missing=not args.no_prune_missing,
        )
    finally:
        if _sigil_client is not None:
            _sigil_client.shutdown()

if __name__ == "__main__":
    main()
